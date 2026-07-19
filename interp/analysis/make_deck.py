"""Generate a shareable overview deck of the spatial-interp project.
Run: interp/.venv/bin/python interp/analysis/make_deck.py
Output: interp/geometry_interp_overview.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# palette
NAVY  = RGBColor(0x1F, 0x2A, 0x44)
BLUE  = RGBColor(0x2E, 0x5C, 0xA8)
GREEN = RGBColor(0x1B, 0x7F, 0x4B)
AMBER = RGBColor(0xB5, 0x6A, 0x00)
GREY  = RGBColor(0x55, 0x5B, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = prs.slide_width


def _box(slide, l, t, w, h):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True
    return tf


def _bar(slide, color, t, h):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(t), W, Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def title_slide(title, subtitle, tagline):
    s = prs.slides.add_slide(BLANK)
    _bar(s, NAVY, 0, 7.5)
    tf = _box(s, 0.9, 2.2, 11.5, 2.2)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(20); r2.font.color.rgb = RGBColor(0xC9, 0xD6, 0xEC)
    p2.space_before = Pt(14)
    tf2 = _box(s, 0.9, 5.6, 11.5, 1.2)
    rp = tf2.paragraphs[0]; rr = rp.add_run(); rr.text = tagline
    rr.font.size = Pt(15); rr.font.italic = True; rr.font.color.rgb = RGBColor(0x9F,0xB2,0xD4)


def content_slide(tag, tag_color, takeaway, bullets, notes=""):
    s = prs.slides.add_slide(BLANK)
    _bar(s, NAVY, 0, 1.35)
    # tag chip
    chip = _box(s, 0.6, 0.18, 3.0, 0.5)
    cp = chip.paragraphs[0]; cr = cp.add_run(); cr.text = tag
    cr.font.size = Pt(13); cr.font.bold = True; cr.font.color.rgb = tag_color
    # takeaway (the headline of the slide)
    tk = _box(s, 0.6, 0.55, 12.1, 0.75)
    tp = tk.paragraphs[0]; tr = tp.add_run(); tr.text = takeaway
    tr.font.size = Pt(24); tr.font.bold = True; tr.font.color.rgb = WHITE
    # bullets
    body = _box(s, 0.7, 1.7, 12.0, 5.4)
    for i, (txt, lvl, color) in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.level = lvl
        r = p.add_run(); r.text = ("•  " if lvl == 0 else "–  ") + txt
        r.font.size = Pt(19 - 2*lvl)
        r.font.color.rgb = color or RGBColor(0x22,0x26,0x2E)
        p.space_after = Pt(7)
    if notes:
        s.notes_slide.notes_text_frame.text = notes


B = lambda t, lvl=0, c=None: (t, lvl, c)  # bullet helper

# 1 ---------------------------------------------------------------
title_slide(
    "Does an LLM build an internal map of geometry?",
    "Mechanistic interpretability of the geometry-diagram generator  ·  Qwen2.5 7B & 32B",
    "What the model represents internally when it constructs a figure — and how we tested it honestly.")

# 2 ---------------------------------------------------------------
content_slide("THE QUESTION", BLUE,
    "When the model draws a figure, does it 'understand' the geometry — and where?",
    [B("An LLM is matrix math on token vectors. When it writes 'M = midpoint of A,B', is it"),
     B("manipulating symbols, or has it built an internal representation of the space?", 1, GREY),
     B("If it has a representation: WHERE in its layers does it form, and is it in a usable form?"),
     B("This is 'mechanistic interpretability' — opening the model up, not just judging its output.", 1, GREY)],
    "Frame for the team: we're not measuring whether the diagrams look right — we're asking what the model is doing internally. The whole project is about reading the model's 'mind' while it works.")

# 3 ---------------------------------------------------------------
content_slide("WHY OUR SETUP IS SPECIAL", BLUE,
    "Our pipeline gives ground truth for free — a rare advantage.",
    [B("Most LLM interpretability has no answer key (what's the 'correct' internal state for an essay?)."),
     B("Our generator compiles every construction to exact geometry (SymPy): coordinates, angles, relations."),
     B("So for every token the model writes, we KNOW the true geometric fact it should encode."),
     B("That answer key is what makes 'decode the geometry from the activations' even possible.", 1, GREY)],
    "This is the single reason the project is feasible here and not on a generic chatbot. We built a render-free grader + ground-truth extractor before touching any activations.")

# 4 ---------------------------------------------------------------
content_slide("THE METHOD", BLUE,
    "Three steps: capture the model's internal state, decode it, then test causally.",
    [B("CAPTURE — run the model, photograph its internal vectors at every layer, for each token."),
     B("PROBE — train a simple linear readout: can we recover a geometric fact from those vectors?"),
     B("PATCH — swap an internal vector between near-identical prompts; does the output change?"),
     B("Probe = 'is the info there?'  Patch = 'does the model actually use it?'", 1, GREY)],
    "Capture and probe are offline/cheap once done. Patch is the causal gold standard. We staged it so each step only mattered if the prior one passed a gate.")

# 5 ---------------------------------------------------------------
content_slide("KEY INTUITION", BLUE,
    "The 'residual stream': a token's vector is a workspace that each layer rewrites.",
    [B("A token starts as just its word; it flows up through ~28 layers, each ADDING an edit."),
     B("ATTENTION is where tokens exchange info — 'M' pulls in 'A' and 'B' to compute their midpoint."),
     B("So a fact that combines tokens (a midpoint, a position) CAN'T exist at layer 0 — it's built with depth."),
     B("We keep the probe LINEAR so 'decodable' = 'in a form the model's own machinery can read'.", 1, GREY)],
    "If they remember one thing: layers = the time-axis of the model's thinking. A fact appearing only in later layers means the model computed it, didn't just copy it from the words.")

# 6 ---------------------------------------------------------------
content_slide("WHAT WE RAN", BLUE,
    "Two model sizes, on rented GPUs, against real exam-geometry prompts.",
    [B("Models: Qwen2.5-7B (full precision) and Qwen2.5-32B (4-bit, to fit the GPU)."),
     B("Prompts: GenExam geometry problems; the model writes constructions, we keep the valid ones."),
     B("Captured internal activations at all layers; trained probes for coordinates, angles, relations."),
     B("Ops: detached GPU runs, monitored every 20 min, crash-proofed for long unattended captures.", 1, GREY)],
    "Hardware was a constant constraint (memory, disk) — we engineered around it. The 32B is 4-bit, which matters for one caveat later (quantization).")

# 7 ---------------------------------------------------------------
content_slide("RESULT 1 — SOLID", GREEN,
    "Scale clearly helps the task: 32B succeeds 2x as often as 7B.",
    [B("Valid-construction rate:  7B = 20%   →   32B = 40%."),
     B("This is a generation-quality number — no probe involved, so it's rock-solid."),
     B("32B's errors are 'deeper' (geometry reference mistakes), not basic schema confusion."),
     B("Takeaway: a bigger model is meaningfully better at producing correct geometry.", 1, GREEN)],
    "Easiest win to communicate. Pure capability. Unaffected by any of the internal-probe caveats.")

# 8 ---------------------------------------------------------------
content_slide("RESULT 2 — SOLID (the headline finding)", GREEN,
    "The model internally represents WHAT each point is — beyond just its name.",
    [B("At the token naming a point, we can linearly decode its geometric ROLE:"),
     B("midpoint / perpendicular / intersection / foot / ... ", 1, GREY),
     B("Accuracy 0.70 (7B) to 0.86 (32B); +0.37 ABOVE a 'guess from the name alone' baseline."),
     B("Holds up under leak-free evaluation, in BOTH model sizes — our most robust claim.", 1, GREEN)],
    "The +0.37 'lift over naming' is crucial: midpoints are often named M, so we control for that. The signal beyond naming is the real representation. This survived our bug fix (next slides).")

# 9 ---------------------------------------------------------------
content_slide("RESULT 3 — SOLID (causal)", GREEN,
    "The model doesn't just store a value — it uses it.",
    [B("Minimal pairs: 'angle = 60°' vs 'angle = 70°', identical otherwise."),
     B("We overwrite the model's internal vector at the angle, mid-computation."),
     B("The model's output flips 60 → 70 — so that representation causally drives behavior."),
     B("Decoding shows info is PRESENT; patching shows it's USED. This is the stronger test.", 1, GREEN)],
    "Patching is a separate method from the probes, so it's untouched by the leakage issue. It's a modest 'uses the stated value' result, but it's genuinely causal.")

# 10 --------------------------------------------------------------
content_slide("WHAT WE WALKED BACK — and why that's good", AMBER,
    "A code review caught a bug; we re-checked and dropped the over-claims.",
    [B("Bug: the probe was accidentally trained & tested on the SAME prompts (different samples) → it memorized."),
     B("After fixing it and re-running honestly:"),
     B("'Coordinate map' (R² 0.49) collapsed to ~0.15 — mostly memorization, not a real claim.", 1, AMBER),
     B("'Angle' decoding became unstable — also not supported.", 1, AMBER),
     B("The relational-role result (Result 2) SURVIVED the fix. We corrected the writeup openly.", 0, GREEN)],
    "Lead with this as a strength, not an embarrassment: an external review found a real flaw, we fixed it, re-ran, and kept only what held. That's how the process is supposed to work.")

# 11 --------------------------------------------------------------
content_slide("THE BOTTLENECK", AMBER,
    "The real limit isn't the method — it's prompt diversity.",
    [B("Every dataset had only 35–50 UNIQUE prompts."),
     B("We generated many samples per prompt — but that inflates counts, not variety."),
     B("Coarse facts (a point's role) generalize from few prompts; precise coordinates do not."),
     B("To settle the 'coordinate map' question we need MANY MORE UNIQUE prompts, not more samples.", 1, AMBER)],
    "This reframes the negative result usefully: it's a data-design lesson, not a dead end. The fix is concrete and cheap.")

# 12 --------------------------------------------------------------
content_slide("WHERE THIS GOES NEXT", BLUE,
    "Concrete next steps, in priority order.",
    [B("1. Capture with many more UNIQUE prompts (widen the geometry tiers / generate variants)."),
     B("2. Re-run the full probe suite leak-free + multi-seed for error bars."),
     B("3. Extend causal patching from 'stated value' to the relational representation."),
     B("4. (Stretch) sparse-autoencoder features for interpretable geometric concepts."),
     B("All code, tests, and methodology are committed and reproducible.", 1, GREY)],
    "The infrastructure is built and verified; next iterations are mostly data + compute, not new engineering.")

# 13 --------------------------------------------------------------
content_slide("SUMMARY", BLUE,
    "What we can say today — honestly.",
    [B("SOLID: bigger model does geometry 2x better (40% vs 20%).", 0, GREEN),
     B("SOLID: the model internally represents each point's geometric ROLE, beyond naming.", 0, GREEN),
     B("SOLID: it causally uses stated values (patching).", 0, GREEN),
     B("OPEN: a precise internal coordinate 'map' — not yet supported; needs more unique prompts.", 0, AMBER),
     B("Built a tested, reproducible interpretability pipeline + an honest, self-correcting process.", 0, GREY)],
    "One-liner for the team: the model demonstrably represents WHAT KIND of geometric object each point is and uses stated values; whether it holds a precise spatial map is still open, bottlenecked by prompt diversity — and our pipeline is ready to answer it.")

out = "interp/geometry_interp_overview.pptx"
prs.save(out)
print(f"saved {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
